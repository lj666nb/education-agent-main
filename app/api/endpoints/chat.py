from fastapi import APIRouter, Depends, HTTPException, status, BackgroundTasks, Query
from fastapi.responses import StreamingResponse, FileResponse
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Literal
from datetime import datetime, timezone, timedelta
from enum import Enum
from sqlalchemy import or_, distinct, func
import os
import httpx
import json
import asyncio
import logging
import re
import uuid
from sqlalchemy.orm import Session, joinedload
from app.api.dependencies import CurrentUser, get_current_user
from app.db.database import get_db
from app.models.chat import ChatSession as ChatSessionModel, ChatMessage as ChatMessageModel, ChatAttachment as ChatAttachmentModel
from app.core.config import settings
from app.services.intent_detector import get_intent_detector

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/chat", tags=["Chat"])

# 中国时区 (UTC+8)
CHINA_TZ = timezone(timedelta(hours=8))

def _fmt_iso(dt: datetime) -> str:
    """序列化 datetime 为中国时区 ISO 字符串。"""
    if dt.tzinfo is None:
        dt = dt.replace(tzinfo=CHINA_TZ)
    return dt.isoformat()


class ModelType(str, Enum):
    DEEPSEEK_V4_FLASH = "deepseek-v4-flash"
    DEEPSEEK_V4_PRO = "deepseek-v4-pro"
    QWEN35_PLUS = "qwen3.5-plus"
    QWEN36_PLUS = "qwen3.6-plus"


MODEL_TO_PROVIDER = {
    ModelType.DEEPSEEK_V4_FLASH: "deepseek",
    ModelType.DEEPSEEK_V4_PRO: "deepseek",
    ModelType.QWEN35_PLUS: "qwen",
    ModelType.QWEN36_PLUS: "qwen",
}


class Message(BaseModel):
    role: Literal["user", "assistant", "system"]
    content: str


class ChatCompletionRequest(BaseModel):
    chat_id: Optional[str] = None
    model: ModelType = ModelType.DEEPSEEK_V4_FLASH
    messages: List[Message]
    stream: bool = True
    temperature: float = 0.7
    max_tokens: Optional[int] = None
    enable_thinking: bool = False
    enable_websearch: bool = False
    project_id: Optional[str] = None
    file_ids: Optional[List[str]] = None


async def build_project_context(db: Session, project_id: str, user_id: str, user_query: str):
    from app.crud import project as project_crud

    project = project_crud.get_project(db, project_id, user_id)
    if not project:
        return None, []

    prompts = project_crud.get_prompts(db, project_id)
    active_prompts = [p for p in prompts if p.is_active]
    prompt_context = ""
    if active_prompts:
        prompt_lines = ["[项目提示词]"]
        for p in active_prompts:
            prompt_lines.append(f"- {p.name}: {p.content}")
        prompt_context = "\n".join(prompt_lines)

    retrieval_results = project_crud.retrieve_relevant_chunks(db, project_id, user_query, top_k=3)
    rag_context = ""
    sources = []
    if retrieval_results:
        rag_lines = ["[参考文档内容]"]
        for i, result in enumerate(retrieval_results, 1):
            if "chunk" in result:
                chunk = result["chunk"]
                score = result.get("score", 0)
            else:
                chunk = result
                score = result.get("combined_score", 0)
            content = chunk.get("content", "") if isinstance(chunk, dict) else str(chunk)
            doc_name = chunk.get("document_name", "未知文档") if isinstance(chunk, dict) else "未知文档"
            rag_lines.append(f"- [{doc_name}] (相关度: {score:.2f}): {content[:200]}...")
            sources.append({
                "document_name": doc_name,
                "content_snippet": content[:150],
                "score": round(float(score), 2) if score else 0
            })
        rag_context = "\n".join(rag_lines)

    if not prompt_context and not rag_context:
        return None, []

    full_context = "你正在一个项目工作区中与用户对话。\n"
    if prompt_context:
        full_context += prompt_context + "\n"
    if rag_context:
        full_context += rag_context + "\n"
    full_context += "\n请根据上述项目提示词和参考文档内容回答用户的问题。如果参考文档中有相关信息，请结合文档内容作答。"

    return full_context, sources


class ChatCompletionResponse(BaseModel):
    chat_id: str
    message: Message
    created_at: datetime


class ChatSession(BaseModel):
    id: str
    title: str
    user_id: str
    model: str
    project_id: Optional[str] = None
    project_name: Optional[str] = None
    created_at: datetime
    updated_at: datetime
    message_count: int = 0


class ChatSessionCreate(BaseModel):
    title: Optional[str] = None
    model: ModelType = ModelType.DEEPSEEK_V4_FLASH
    project_id: Optional[str] = None


class ChatSessionUpdate(BaseModel):
    title: Optional[str] = None


class ChatHistoryItem(BaseModel):
    id: str
    role: str
    content: str
    created_at: datetime


class SaveMessageRequest(BaseModel):
    chat_id: str
    role: str
    content: str
    reasoning_content: Optional[str] = None
    citations: Optional[list] = None


class SaveMessageResponse(BaseModel):
    id: str
    chat_id: str
    role: str
    content: str
    created_at: datetime


class AttachmentCreate(BaseModel):
    session_id: str
    file_id: str
    file_name: str
    file_type: str


class AttachmentResponse(BaseModel):
    id: str
    session_id: str
    file_id: str
    file_name: str
    file_type: str
    created_at: datetime

    class Config:
        from_attributes = True


MULTIMODAL_MODELS = {ModelType.QWEN35_PLUS, ModelType.QWEN36_PLUS}

AVAILABLE_MODELS = {
    ModelType.DEEPSEEK_V4_FLASH: {
        "name": "DeepSeek V4 Flash",
        "model_id": "deepseek-v4-flash",
        "base_url": settings.DEEPSEEK_BASE_URL,
        "api_key": settings.DEEPSEEK_API_KEY,
        "supports_streaming": True,
        "supports_thinking": True,
    },
    ModelType.DEEPSEEK_V4_PRO: {
        "name": "DeepSeek V4 Pro",
        "model_id": "deepseek-v4-pro",
        "base_url": settings.DEEPSEEK_BASE_URL,
        "api_key": settings.DEEPSEEK_API_KEY,
        "supports_streaming": True,
        "supports_thinking": True,
    },
    ModelType.QWEN35_PLUS: {
        "name": "Qwen3.5 Plus",
        "model_id": "qwen3.5-plus",
        "base_url": settings.QWEN_BASE_URL,
        "api_key": settings.QWEN_API_KEY,
        "supports_streaming": True,
        "supports_thinking": False,
    },
    ModelType.QWEN36_PLUS: {
        "name": "Qwen3.6 Plus",
        "model_id": "qwen3.6-plus",
        "base_url": settings.QWEN_BASE_URL,
        "api_key": settings.QWEN_API_KEY,
        "supports_streaming": True,
        "supports_thinking": False,
    },
}


def get_model_config(model: ModelType) -> Dict:
    if model not in AVAILABLE_MODELS:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"不支持的模型: {model}"
        )
    return AVAILABLE_MODELS[model]


async def call_llm_api(
    messages: List[Dict[str, str]],
    model: ModelType,
    stream: bool = True,
    temperature: float = 0.7,
    max_tokens: Optional[int] = None,
    enable_thinking: bool = False,
    api_key_override: Optional[str] = None,
    base_url_override: Optional[str] = None,
):
    model_config = get_model_config(model)

    if api_key_override:
        model_config = {**model_config, "api_key": api_key_override}
    if base_url_override:
        model_config = {**model_config, "base_url": base_url_override}

    if not model_config["api_key"]:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"{model.value} API Key 未配置"
        )

    headers = {
        "Authorization": f"Bearer {model_config['api_key']}",
        "Content-Type": "application/json",
    }

    payload = {
        "model": model_config["model_id"],
        "messages": messages,
        "temperature": temperature,
        "stream": stream,
    }

    if max_tokens:
        payload["max_tokens"] = max_tokens

    if enable_thinking:
        if model in [ModelType.DEEPSEEK_V4_FLASH, ModelType.DEEPSEEK_V4_PRO]:
            payload["thinking"] = {"type": "enabled"}
        elif model in [ModelType.QWEN35_PLUS, ModelType.QWEN36_PLUS]:
            payload["extra_body"] = {"enable_thinking": True}

    async with httpx.AsyncClient(timeout=120.0) as client:
        response = await client.post(
            f"{model_config['base_url']}/chat/completions",
            headers=headers,
            json=payload
        )

        if response.status_code != 200:
            raise HTTPException(
                status_code=status.HTTP_502_BAD_GATEWAY,
                detail=f"{model.value} API 调用失败: {response.text}"
            )

        if stream:
            return response
        else:
            return response.json()


async def generate_title_from_first_message(message: str, model: ModelType) -> str:
    prompt = f"请为以下对话生成一个简短的中文标题（不超过20字），仅返回标题不要其他内容：\n\n用户：{message[:50]}..."

    try:
        result = await call_llm_api(
            messages=[{"role": "user", "content": prompt}],
            model=model,
            stream=False,
            temperature=0.3
        )
        title = result["choices"][0]["message"]["content"].strip()
        if len(title) > 20:
            title = title[:20]
        return title or "新对话"
    except Exception:
        return "新对话"


def _is_garbled_text(text: str) -> bool:
    """检测PDF提取的文本是否为乱码（CJK编码错误导致的无意义字符）

    有效的中文PDF文本应包含一定比例的CJK统一表意文字。
    乱码文本通常包含大量无效Unicode字符或替换字符。
    """
    if not text or len(text.strip()) < 50:
        return True  # 太短，无法判断
    total = len(text)
    # 替换字符是乱码的强信号
    if '�' in text:
        return True
    # 统计各类字符
    cjk_count = 0
    ascii_count = 0
    odd_count = 0  # 非CJK、非ASCII、非标点的字符（乱码通常聚集在此）
    for ch in text:
        cp = ord(ch)
        if 0x4E00 <= cp <= 0x9FFF:  # CJK统一表意文字
            cjk_count += 1
        elif 0x3400 <= cp <= 0x4DBF:  # CJK扩展A
            cjk_count += 1
        elif 0x20000 <= cp <= 0x2A6DF:  # CJK扩展B
            cjk_count += 1
        elif 0x20 <= cp <= 0x7E:  # ASCII可打印字符
            ascii_count += 1
        elif cp in (0x0A, 0x0D, 0x09, 0x0C):  # 空白字符
            pass
        elif 0x3000 <= cp <= 0x303F:  # CJK符号和标点
            pass
        elif 0xFF00 <= cp <= 0xFFEF:  # 全角字符
            pass
        elif 0x2000 <= cp <= 0x206F:  # 通用标点
            pass
        elif 0x00A0 <= cp <= 0x00FF:  # Latin-1 Supplement（UTF-8中文解码为Latin-1时的典型范围）
            odd_count += 1
        else:
            odd_count += 1
    # 长文本没有CJK字符 → 乱码（除非绝大多数是ASCII，如代码）
    if total >= 100 and cjk_count == 0 and odd_count > ascii_count * 0.3:
        return True
    # 中文PDF中CJK比例极低 → 乱码
    if total >= 500 and cjk_count / total < 0.01:
        return True
    return False


async def load_and_process_file(db: Session, file_id: str, student_id: str, is_multimodal: bool = False) -> Optional[Dict]:
    import base64
    from app.core.ocr import get_ocr_service

    for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.pdf', '.pptx', '.ppt', '.docx', '.doc']:
        filepath = os.path.join("uploads", f"{file_id}{ext}")
        if os.path.exists(filepath):
            try:
                with open(filepath, 'rb') as f:
                    content = f.read()

                is_image = ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']
                is_pptx = ext == '.pptx'
                is_ppt = ext == '.ppt'
                is_docx = ext == '.docx'
                is_doc = ext == '.doc'
                mime_type = f"image/{ext[1:]}" if is_image else "application/pdf"
                image_base64 = base64.b64encode(content).decode('utf-8')

                result = {
                    "type": "image" if is_image else "pdf",
                    "base64": image_base64,
                    "mime_type": mime_type,
                    "text": None
                }

                if is_ppt:
                    result["type"] = "document"
                    result["text"] = "[旧版PPT格式（.ppt），无法直接解析文本内容。建议将文件另存为 .pptx 格式后重新上传。]"
                    return result

                if is_doc:
                    result["type"] = "document"
                    result["text"] = "[旧版Word格式（.doc），无法直接解析文本内容。建议将文件另存为 .docx 格式后重新上传。]"
                    return result

                if is_pptx:
                    try:
                        from pptx import Presentation
                        from io import BytesIO
                        prs = Presentation(BytesIO(content))
                        slide_texts = []
                        for slide_num, slide in enumerate(prs.slides, 1):
                            texts = []
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text.strip():
                                    texts.append(shape.text.strip())
                                if shape.has_table:
                                    table = shape.table
                                    for row in table.rows:
                                        row_texts = [cell.text.strip() for cell in row.cells]
                                        texts.append(" | ".join(row_texts))
                            if texts:
                                slide_texts.append(f"[第{slide_num}页] " + "\n".join(texts))
                        doc_text = "\n\n".join(slide_texts) if slide_texts else "[PPT文件为空]"
                        if len(doc_text) > 20000:
                            doc_text = doc_text[:20000] + "\n\n... [内容已截断，仅显示前20000字符]"
                        result["type"] = "document"
                        result["text"] = f"[PPT 文件内容]:\n{doc_text}"
                    except ImportError:
                        result["text"] = "[PPT 解析库未安装，无法提取文件内容]"
                    except Exception as e:
                        result["text"] = f"[PPT 文件解析失败: {str(e)}]"
                    return result

                if is_docx:
                    try:
                        from app.services.docx_parser import extract_docx_text
                        doc_text = extract_docx_text(content, max_chars=20000)
                        if not doc_text.startswith("["):
                            doc_text = f"[Word 文件内容]:\n{doc_text}"
                        result["type"] = "document"
                        result["text"] = doc_text
                    except ImportError:
                        result["text"] = "[Word 解析库未安装，无法提取文件内容]"
                    except Exception as e:
                        result["text"] = f"[Word 文件解析失败: {str(e)}]"
                    return result

                if is_image and not is_multimodal:
                    ocr_service = get_ocr_service(student_id)
                    if ocr_service:
                        try:
                            ocr_result = ocr_service.recognize_text(image_base64)
                            texts = []
                            if 'words_result' in ocr_result:
                                texts = [item.get('words', '') for item in ocr_result['words_result']]
                            if texts:
                                result["text"] = f"[图片 OCR 识别结果]:\n" + "\n".join(texts)
                        except:
                            pass
                    if not result["text"]:
                        result["text"] = f"[图片文件，内容无法识别]"
                elif is_image and is_multimodal:
                    result["text"] = f"[图片文件]"
                else:
                    pdf_text = None
                    total_pages = 0
                    page_images = []
                    try:
                        import fitz
                        doc = fitz.open(stream=content, filetype="pdf")
                        total_pages = len(doc)
                        page_texts = []

                        for page_num in range(total_pages):
                            page = doc[page_num]
                            page_text = page.get_text().strip()
                            if page_text:
                                page_texts.append(f"[第{page_num + 1}页] {page_text}")

                        total_extracted = sum(len(t) for t in page_texts)
                        has_meaningful_text = total_extracted > 100

                        if has_meaningful_text:
                            full_text = "\n\n".join(page_texts)
                            # 检查是否为乱码（编码错误导致的无意义字符）
                            if _is_garbled_text(full_text):
                                logger.warning(
                                    f"PDF {file_id}: extracted {total_extracted} chars but text is garbled, "
                                    f"falling back to image rendering"
                                )
                                has_meaningful_text = False

                        if has_meaningful_text:
                            # ── 文本有效，直接使用 ──
                            if len(full_text) > 20000:
                                full_text = full_text[:20000] + "\n\n... [内容已截断，仅显示前20000字符]"
                            pdf_text = full_text
                        else:
                            # ── 文本无效（太少或乱码），降级为图片渲染 ──
                            max_pages_for_images = 15
                            pages_to_render = min(total_pages, max_pages_for_images)

                            if is_multimodal:
                                # 多模态模型：渲染页面为图片，直接发给模型阅读
                                render_dpi = 200
                                for page_num in range(pages_to_render):
                                    page = doc[page_num]
                                    pix = page.get_pixmap(dpi=render_dpi)
                                    img_bytes = pix.tobytes("jpeg", jpg_quality=85)
                                    img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                                    page_images.append({
                                        "base64": img_b64,
                                        "mime_type": "image/jpeg",
                                        "page_num": page_num + 1
                                    })
                                total_shown = len(page_images)
                                extra_note = ""
                                if total_pages > total_shown:
                                    extra_note = f"（PDF共{total_pages}页，仅展示前{total_shown}页）"
                                pdf_text = f"以下是该PDF文件第1-{total_shown}页的内容截图，请直接阅读图片。{extra_note}"
                            else:
                                # 非多模态模型：尝试OCR，无OCR时给出引导
                                ocr_service = get_ocr_service(student_id)
                                ocr_texts = []
                                for page_num in range(pages_to_render):
                                    page = doc[page_num]
                                    pix = page.get_pixmap(dpi=200)
                                    img_bytes = pix.tobytes("jpeg", jpg_quality=85)
                                    img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                                    if ocr_service:
                                        try:
                                            ocr_result = ocr_service.recognize_text(img_b64)
                                            words = []
                                            if 'words_result' in ocr_result:
                                                words = [item.get('words', '') for item in ocr_result['words_result']]
                                            if words:
                                                ocr_texts.append(f"[第{page_num + 1}页] " + "\n".join(words))
                                        except:
                                            pass
                                if ocr_texts:
                                    pdf_text = "\n\n".join(ocr_texts)
                                    if len(pdf_text) > 20000:
                                        pdf_text = pdf_text[:20000] + "\n\n... [内容已截断，仅显示前20000字符]"
                                    if total_pages > pages_to_render:
                                        pdf_text += f"\n\n[注意：PDF共{total_pages}页，仅OCR识别了前{pages_to_render}页]"
                                else:
                                    if ocr_service:
                                        pdf_text = (
                                            f"[该PDF共{total_pages}页，尝试使用OCR识别但未能提取到有效文字。\n"
                                            f"可能原因：PDF内容为特殊格式或图片质量较低。\n"
                                            f"建议切换到 Qwen3.5-plus 或 Qwen3.6-plus 模型（支持直接阅读PDF图片）。]"
                                        )
                                    else:
                                        pdf_text = (
                                            f"[该PDF共{total_pages}页，无法提取到有效文字内容。\n"
                                            f"可能原因：该PDF为扫描件或图片型PDF，或文字编码无法解析。\n"
                                            f"解决方法：\n"
                                            f"1. 在「设置」→「API 配置」中配置百度OCR（推荐）\n"
                                            f"2. 切换到 Qwen3.5-plus 或 Qwen3.6-plus 模型（支持直接阅读PDF图片）\n"
                                            f"3. 上传文字型PDF（从Word/LaTeX直接导出）]"
                                        )
                                # 非多模态模型也渲染前3页图片供参考（部分API如DeepSeek v4支持图片输入）
                                extra_images = min(total_pages, 3)
                                for page_num in range(extra_images):
                                    page = doc[page_num]
                                    pix = page.get_pixmap(dpi=200)
                                    img_bytes = pix.tobytes("jpeg", jpg_quality=85)
                                    img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                                    page_images.append({
                                        "base64": img_b64,
                                        "mime_type": "image/jpeg",
                                        "page_num": page_num + 1
                                    })
                        doc.close()
                    except ImportError:
                        pdf_text = "[PDF 解析库未安装，无法提取文件内容]"
                    except Exception as e:
                        pdf_text = f"[PDF 文件解析失败: {str(e)}]"

                    if page_images:
                        result["type"] = "pdf_with_images"
                        result["page_images"] = page_images
                    if pdf_text:
                        result["text"] = f"[PDF 文件内容]:\n{pdf_text}"
                    else:
                        result["text"] = "[PDF 文件，无法提取任何内容]"

                return result
            except Exception as e:
                return {"type": "error", "text": f"[文件读取失败: {str(e)}]", "base64": None, "mime_type": None}
    return None


@router.post("/completions")
async def chat_completions(
    request: ChatCompletionRequest,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    from app.crud.api_settings import api_settings_crud

    if not request.messages:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="消息列表不能为空"
        )

    provider = MODEL_TO_PROVIDER.get(request.model)
    if provider:
        is_available = api_settings_crud.is_provider_available(db, str(current_user.student_id), provider)
        if not is_available:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail=f"当前模型 {request.model.value} 不可用，请先在设置中配置 {provider} API"
            )

    messages_dict = [{"role": msg.role, "content": msg.content} for msg in request.messages]

    is_multimodal = request.model in MULTIMODAL_MODELS

    if request.file_ids:
        file_contents = []
        for file_id in request.file_ids:
            file_content = await load_and_process_file(db, file_id, str(current_user.student_id), is_multimodal)
            if file_content:
                file_contents.append(file_content)

        if file_contents:
            for msg_dict in reversed(messages_dict):
                if msg_dict["role"] == "user":
                    if is_multimodal:
                        if isinstance(msg_dict["content"], str):
                            msg_dict["content"] = [{"type": "text", "text": msg_dict["content"]}]
                        for fc in file_contents:
                            if fc.get("type") == "pdf_with_images" and fc.get("page_images"):
                                for page_img in fc["page_images"]:
                                    msg_dict["content"].append({
                                        "type": "image_url",
                                        "image_url": {"url": f"data:{page_img['mime_type']};base64,{page_img['base64']}"}
                                    })
                                if fc.get("text"):
                                    msg_dict["content"].append({"type": "text", "text": f"\n[附件说明]: {fc['text']}"})
                            elif fc.get("type") == "image" and fc.get("base64"):
                                msg_dict["content"].append({
                                    "type": "image_url",
                                    "image_url": {"url": f"data:{fc['mime_type']};base64,{fc['base64']}"}
                                })
                                if fc.get("text"):
                                    msg_dict["content"].append({"type": "text", "text": f"\n[附件内容]: {fc['text']}"})
                            elif fc.get("text"):
                                msg_dict["content"].append({"type": "text", "text": f"\n[附件内容]: {fc['text']}"})
                    else:
                        # 非多模态模型：优先使用文本，如有 page_images 也一并发送
                        # （DeepSeek v4 等模型也支持图片输入）
                        has_images = any(
                            fc.get("type") == "pdf_with_images" and fc.get("page_images")
                            for fc in file_contents
                        )
                        if has_images:
                            if isinstance(msg_dict["content"], str):
                                msg_dict["content"] = [{"type": "text", "text": msg_dict["content"]}]
                            for fc in file_contents:
                                if fc.get("type") == "pdf_with_images" and fc.get("page_images"):
                                    for page_img in fc["page_images"]:
                                        msg_dict["content"].append({
                                            "type": "image_url",
                                            "image_url": {"url": f"data:{page_img['mime_type']};base64,{page_img['base64']}"}
                                        })
                                    if fc.get("text"):
                                        msg_dict["content"].append({"type": "text", "text": f"\n[附件说明]: {fc['text']}"})
                                elif fc.get("text"):
                                    msg_dict["content"].append({"type": "text", "text": f"\n[附件内容]: {fc['text']}"})
                        else:
                            file_info = "\n\n[用户上传的文件内容]\n" + "\n---\n".join([fc.get("text", "") for fc in file_contents])
                            msg_dict["content"] = msg_dict["content"] + file_info
                    break

    web_citations = None
    search_error = None
    if request.enable_websearch:
        from app.core.web_search import perform_search_and_build_context, CITATION_SYSTEM_PROMPT
        user_query = ""
        for msg in reversed(messages_dict):
            if msg["role"] == "user":
                user_query = msg["content"] if isinstance(msg["content"], str) else str(msg["content"])
                break
        if user_query:
            try:
                # 优先使用用户配置的 Tavily API Key
                user_tavily_key = None
                try:
                    tavily_setting = api_settings_crud.get_setting_value(db, str(current_user.student_id), "tavily")
                    if tavily_setting:
                        user_tavily_key = tavily_setting.get("api_key")
                except Exception:
                    pass
                citations, context = await perform_search_and_build_context(user_query, api_key=user_tavily_key)
                if citations and context:
                    web_citations = citations
                    # 找到最后一个 user 消息的索引，在其前插入 system 消息
                    user_idx = None
                    for idx in range(len(messages_dict) - 1, -1, -1):
                        if messages_dict[idx]["role"] == "user":
                            user_idx = idx
                            break
                    if user_idx is not None:
                        system_msg = {
                            "role": "system",
                            "content": CITATION_SYSTEM_PROMPT + "\n\n" + context,
                        }
                        messages_dict.insert(user_idx, system_msg)
                elif citations is not None and not citations:
                    search_error = "搜索未返回结果"
            except Exception as e:
                logger.warning(f"联网搜索异常: {e}")
                search_error = str(e)

    if request.project_id:
        user_query = ""
        for msg in reversed(request.messages):
            if msg.role == "user":
                user_query = msg.content
                break
        rag_sources = []
        if user_query:
            project_context, rag_sources = await build_project_context(db, request.project_id, str(current_user.student_id), user_query)
            if project_context:
                messages_dict.insert(0, {"role": "system", "content": project_context})
    else:
        rag_sources = []

    # 意图检测（获取用户最后一条消息）
    intent_result = None
    user_last_message = ""
    for msg in reversed(request.messages):
        if msg.role == "user":
            user_last_message = msg.content
            break
    if user_last_message:
        try:
            # 获取用户自己的 API Key 用于意图检测
            user_api_key = None
            user_api_base = None
            user_api_model = None
            provider = MODEL_TO_PROVIDER.get(request.model)
            if provider:
                user_api = api_settings_crud.get_setting_value(db, str(current_user.student_id), provider)
                if user_api:
                    user_api_key = user_api.get("api_key")
                    user_api_base = user_api.get("base_url")
                    user_api_model = get_model_config(request.model)["model_id"]
            if user_api_key:
                from app.services.intent_detector import IntentDetector
                detector = IntentDetector(api_key=user_api_key, base_url=user_api_base, model=user_api_model)
            else:
                detector = get_intent_detector()
            intent_result = await detector.analyze(user_last_message)
            if intent_result and intent_result.get("knowledge_points"):
                logger.info(f"意图检测: 学习相关, 知识点={intent_result['knowledge_points']}")
            elif intent_result and not intent_result.get("is_learning_related", True):
                logger.info(f"意图检测: 无关内容, confidence={intent_result.get('confidence')}")
        except Exception as e:
            logger.warning(f"意图检测异常: {e}")

    try:
        if request.stream:
            async def stream_generator():
                async with httpx.AsyncClient(timeout=120.0) as client:
                    model_config = get_model_config(request.model)
                    provider = MODEL_TO_PROVIDER.get(request.model)
                    if provider:
                        user_api = api_settings_crud.get_setting_value(db, str(current_user.student_id), provider)
                        if user_api:
                            model_config = {**model_config, "api_key": user_api["api_key"]}
                            if user_api.get("base_url"):
                                model_config["base_url"] = user_api["base_url"]

                    headers = {
                        "Authorization": f"Bearer {model_config['api_key']}",
                        "Content-Type": "application/json",
                    }

                    payload = {
                        "model": model_config["model_id"],
                        "messages": messages_dict,
                        "temperature": request.temperature,
                        "stream": True,
                    }

                    if request.max_tokens:
                        payload["max_tokens"] = request.max_tokens

                    if request.enable_thinking and request.model in [ModelType.DEEPSEEK_V4_FLASH, ModelType.DEEPSEEK_V4_PRO]:
                        payload["thinking"] = {"type": "enabled"}
                    elif request.enable_thinking and request.model in [ModelType.QWEN35_PLUS, ModelType.QWEN36_PLUS]:
                        payload["extra_body"] = {"enable_thinking": True}

                    chat_id = request.chat_id or f"chat_{current_user.student_id}_{datetime.now().timestamp()}"

                    # 1. 如果检测到无关内容，先在流开头发出警告事件
                    if intent_result and not intent_result.get("is_learning_related", True):
                        irrelevant_event = {
                            'type': 'irrelevant_content',
                            'message': '当前为学习助手，请提问学科相关问题',
                            'confidence': intent_result.get('confidence'),
                        }
                        yield f"data: {json.dumps(irrelevant_event)}\n\n"

                    # 2. 联网搜索状态事件
                    if web_citations is not None:
                        yield f"data: {json.dumps({'type': 'search_status', 'status': 'searching', 'message': '正在搜索...'})}\n\n"
                        yield f"data: {json.dumps({'type': 'search_status', 'status': 'done', 'count': len(web_citations), 'message': f'已搜索到 {len(web_citations)} 个网页'})}\n\n"
                    elif search_error:
                        yield f"data: {json.dumps({'type': 'search_status', 'status': 'searching', 'message': '正在搜索...'})}\n\n"
                        yield f"data: {json.dumps({'type': 'search_status', 'status': 'error', 'message': f'搜索暂时不可用: {search_error}'})}\n\n"

                    try:
                        async with client.stream(
                            "POST",
                            f"{model_config['base_url']}/chat/completions",
                            headers=headers,
                            json=payload
                        ) as response:
                            if response.status_code != 200:
                                # StreamingResponse has already sent HTTP 200 headers at this
                                # point, so raising here only aborts the connection and the
                                # browser reports a misleading "network error". Read the
                                # upstream response and return a structured SSE error instead.
                                await response.aread()
                                upstream_detail = response.text.strip()[:500]
                                if response.status_code == 401:
                                    message = "AI 服务鉴权失败：API Key 无效或已过期，请在 API 设置中重新配置"
                                elif response.status_code == 403:
                                    message = "AI 服务拒绝访问：当前 API Key 没有该模型权限"
                                elif response.status_code == 429:
                                    message = "AI 服务请求过于频繁或余额不足，请稍后重试"
                                else:
                                    message = f"{request.model.value} API 调用失败（{response.status_code}）"
                                    if upstream_detail:
                                        message += f"：{upstream_detail}"
                                logger.warning(
                                    "AI 上游调用失败: model=%s status=%s detail=%s",
                                    request.model.value,
                                    response.status_code,
                                    upstream_detail,
                                )
                                yield f"data: {json.dumps({'type': 'error', 'status': response.status_code, 'message': message})}\n\n"
                                return

                            async for line in response.aiter_lines():
                                if line and line.startswith("data: "):
                                    data = line[6:]
                                    if data == "[DONE]":
                                        yield f"data: {json.dumps({'done': True})}\n\n"
                                    else:
                                        try:
                                            chunk = json.loads(data)
                                            if "choices" in chunk and len(chunk["choices"]) > 0:
                                                delta = chunk["choices"][0].get("delta", {})
                                                content = delta.get("content", "")
                                                reasoning_content = delta.get("reasoning_content", "")
                                                if content or reasoning_content:
                                                    resp_data = {
                                                        'content': content,
                                                        'reasoning_content': reasoning_content if reasoning_content else None,
                                                        'chat_id': chat_id
                                                    }
                                                    yield f"data: {json.dumps(resp_data)}\n\n"
                                        except json.JSONDecodeError:
                                            continue
                    except httpx.TimeoutException:
                        logger.warning("AI 上游调用超时: model=%s", request.model.value)
                        yield f"data: {json.dumps({'type': 'error', 'message': 'AI 服务响应超时，请稍后重试'})}\n\n"
                        return
                    except httpx.RequestError as exc:
                        logger.warning(
                            "AI 上游网络请求失败: model=%s error=%s",
                            request.model.value,
                            type(exc).__name__,
                        )
                        yield f"data: {json.dumps({'type': 'error', 'message': '服务器暂时无法连接 AI 服务，请稍后重试'})}\n\n"
                        return

                    # 4. 引用数据（流结束后发送）
                    if web_citations:
                        yield f"data: {json.dumps({'type': 'citations', 'citations': web_citations})}\n\n"

                    final_data = {'done': True, 'thinking_done': True}
                    if rag_sources:
                        final_data['sources'] = rag_sources
                    if web_citations:
                        final_data['citations'] = web_citations

                    yield f"data: {json.dumps(final_data)}\n\n"

            return StreamingResponse(
                stream_generator(),
                media_type="text/event-stream",
                headers={
                    "Cache-Control": "no-cache",
                    "Connection": "keep-alive",
                    "X-Accel-Buffering": "no",
                }
            )
        else:
            user_api_key = None
            user_api_base_url = None
            provider = MODEL_TO_PROVIDER.get(request.model)
            if provider:
                user_api = api_settings_crud.get_setting_value(db, str(current_user.student_id), provider)
                if user_api:
                    user_api_key = user_api["api_key"]
                    user_api_base_url = user_api.get("base_url")

            result = await call_llm_api(
                messages=messages_dict,
                model=request.model,
                stream=False,
                temperature=request.temperature,
                max_tokens=request.max_tokens,
                api_key_override=user_api_key,
                base_url_override=user_api_base_url,
            )

            chat_id = request.chat_id or f"chat_{current_user.student_id}_{datetime.now().timestamp()}"
            content = result["choices"][0]["message"]["content"]

            result_obj = {
                "chat_id": chat_id,
                "message": {
                    "role": "assistant",
                    "content": content
                },
                "created_at": _fmt_iso(datetime.now())
            }
            if web_citations:
                result_obj["citations"] = web_citations
            return result_obj

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"对话生成失败: {str(e)}"
        )


class IntentAnalysisRequest(BaseModel):
    message: str


class IntentAnalysisResponse(BaseModel):
    is_learning_related: bool
    confidence: float
    reason: str
    knowledge_points: List[dict] = []


@router.post("/intent-analyze")
async def analyze_intent(
    request: IntentAnalysisRequest,
    current_user: CurrentUser = Depends(get_current_user),
):
    """分析用户消息意图和知识盲区（不触发资源生成）"""
    detector = get_intent_detector()
    result = await detector.analyze(request.message)
    return IntentAnalysisResponse(
        is_learning_related=result.get("is_learning_related", True),
        confidence=result.get("confidence", 0.5),
        reason=result.get("reason", ""),
        knowledge_points=result.get("knowledge_points", []),
    )


@router.get("/models")
async def list_models(
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    from app.crud.api_settings import api_settings_crud

    models_data = []
    for model_type, config in AVAILABLE_MODELS.items():
        provider = MODEL_TO_PROVIDER.get(model_type)
        is_available = False
        if provider:
            is_available = api_settings_crud.is_provider_available(db, str(current_user.student_id), provider)
        models_data.append({
            "id": model_type.value,
            "name": config["name"],
            "supports_streaming": config["supports_streaming"],
            "supports_thinking": config["supports_thinking"],
            "is_available": is_available,
        })
    return {"models": models_data}


@router.get("/history")
async def get_chat_history(
    limit: int = 50,
    offset: int = 0,
    search: Optional[str] = Query(None, description="搜索关键词，会搜索标题和消息内容"),
    project_id: Optional[str] = Query(None, description="按项目ID筛选"),
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    base_query = db.query(ChatSessionModel).filter(
        ChatSessionModel.user_id == current_user.student_id
    )

    if project_id is None or project_id == "" or project_id == "null":
        base_query = base_query.filter(ChatSessionModel.project_id == None)
    elif project_id:
        base_query = base_query.filter(ChatSessionModel.project_id == project_id)

    if search:
        search_pattern = f"%{search}%"
        base_query = base_query.join(ChatMessageModel, ChatSessionModel.id == ChatMessageModel.session_id).filter(
            or_(
                ChatSessionModel.title.ilike(search_pattern),
                ChatMessageModel.content.ilike(search_pattern)
            )
        ).distinct()

    sessions = base_query.order_by(ChatSessionModel.updated_at.desc()).offset(offset).limit(limit).all()

    count_query = db.query(func.count(distinct(ChatSessionModel.id))).filter(
        ChatSessionModel.user_id == current_user.student_id
    )
    if project_id is None or project_id == "" or project_id == "null":
        count_query = count_query.filter(ChatSessionModel.project_id == None)
    elif project_id:
        count_query = count_query.filter(ChatSessionModel.project_id == project_id)
    if search:
        search_pattern = f"%{search}%"
        count_query = count_query.join(ChatMessageModel, ChatSessionModel.id == ChatMessageModel.session_id).filter(
            or_(
                ChatSessionModel.title.ilike(search_pattern),
                ChatMessageModel.content.ilike(search_pattern)
            )
        )
    total = count_query.scalar()

    return {
        "chats": [
            {
                "id": session.id,
                "title": session.title,
                "user_id": str(session.user_id),
                "project_id": session.project_id,
                "project_name": session.project.name if session.project else None,
                "model": session.model,
                "created_at": _fmt_iso(session.created_at),
                "updated_at": _fmt_iso(session.updated_at),
                "message_count": session.message_count
            }
            for session in sessions
        ],
        "total": total,
        "limit": limit,
        "offset": offset,
        "search": search
    }


@router.get("/{chat_id}/messages")
async def get_chat_messages(
    chat_id: str,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    messages = db.query(ChatMessageModel).filter(
        ChatMessageModel.session_id == chat_id
    ).order_by(ChatMessageModel.created_at).all()

    return {
        "chat_id": chat_id,
        "messages": [
            {
                "id": str(msg.id),
                "role": msg.role,
                "content": msg.content,
                "reasoning_content": msg.reasoning_content,
                "citations": msg.citations,
                "created_at": _fmt_iso(msg.created_at)
            }
            for msg in messages
        ]
    }


@router.post("/sessions")
async def create_chat_session(
    request: ChatSessionCreate,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    chat_id = f"chat_{current_user.student_id}_{int(datetime.now().timestamp())}"

    project_name = None
    if request.project_id:
        from app.crud import project as project_crud
        project = project_crud.get_project(db, request.project_id, str(current_user.student_id))
        if project:
            project_name = project.name

    db_session = ChatSessionModel(
        id=chat_id,
        user_id=current_user.student_id,
        project_id=request.project_id,
        title=request.title or "新对话",
        model=request.model.value,
        message_count=0
    )
    db.add(db_session)
    db.commit()
    db.refresh(db_session)

    return {
        "id": db_session.id,
        "title": db_session.title,
        "user_id": str(db_session.user_id),
        "project_id": db_session.project_id,
        "project_name": project_name,
        "model": db_session.model,
        "created_at": _fmt_iso(db_session.created_at),
        "updated_at": _fmt_iso(db_session.updated_at),
        "message_count": db_session.message_count
    }


@router.patch("/sessions/{chat_id}")
async def update_chat_session(
    chat_id: str,
    request: ChatSessionUpdate,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    db_session = db.query(ChatSessionModel).filter(ChatSessionModel.id == chat_id).first()
    if db_session:
        if request.title:
            db_session.title = request.title
        db.commit()
        db.refresh(db_session)

    return {
        "id": chat_id,
        "title": db_session.title if db_session else (request.title or "新对话"),
        "updated_at": _fmt_iso(db_session.updated_at) if db_session else _fmt_iso(datetime.now())
    }


@router.delete("/sessions/{chat_id}")
async def delete_chat_session(
    chat_id: str,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    db_session = db.query(ChatSessionModel).filter(ChatSessionModel.id == chat_id).first()
    if db_session:
        db.delete(db_session)
        db.commit()
    return {"success": True, "message": "对话已删除"}


@router.post("/messages")
async def save_message(
    request: SaveMessageRequest,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    content = request.content

    # ═══ 处理 [PLOT] 代码块：执行 → 保存 PNG → 替换为图片 URL ═══
    if request.role == "assistant" and "[PLOT]" in content:
        def _process_plot_block(match):
            """执行 [PLOT] 代码，返回 markdown 图片语法"""
            code = match.group(1).strip()
            code = re.sub(r'^```[a-zA-Z]*\n?', '', code)
            code = re.sub(r'\n?```$', '', code)
            try:
                from app.api.endpoints.code_execution import _execute_plot
                result = _execute_plot(code)
                if result.success and result.image:
                    # 提取 base64 数据并保存为文件
                    b64_data = result.image
                    if b64_data.startswith("data:image/png;base64,"):
                        b64_data = b64_data[len("data:image/png;base64,"):]
                    import base64
                    img_bytes = base64.b64decode(b64_data)

                    plot_id = uuid.uuid4().hex[:12]
                    filename = f"plot_{plot_id}.png"
                    plot_dir = os.path.join("uploads", "plots")
                    os.makedirs(plot_dir, exist_ok=True)
                    filepath = os.path.join(plot_dir, filename)
                    with open(filepath, "wb") as f:
                        f.write(img_bytes)

                    url = f"/api/v1/chat/plots/{filename}"
                    logger.info(f"[PLOT] 已保存: {filepath} ({len(img_bytes)} bytes)")
                    return f"\n\n![图表]({url})\n\n"
                else:
                    logger.warning(f"[PLOT] 执行失败: {result.stderr}")
                    return f"\n\n> ⚠️ 图表生成失败: {result.stderr}\n\n"
            except Exception as e:
                logger.error(f"[PLOT] 处理异常: {e}")
                return f"\n\n> ⚠️ 图表处理异常: {str(e)}\n\n"

        content = re.sub(r'\[PLOT\]([\s\S]*?)\[/PLOT\]', _process_plot_block, content)

    # ═══ 处理 [DRAWIO] 代码块：尝试渲染为 PNG → 替换为图片 URL ═══
    if request.role == "assistant" and "[DRAWIO]" in content:
        def _process_drawio_block(match):
            xml = match.group(1).strip()
            xml = re.sub(r'^```[a-zA-Z]*\n?', '', xml)
            xml = re.sub(r'\n?```$', '', xml)
            if not xml.strip():
                return "\n\n> ⚠️ 图表内容为空\n\n"
            try:
                from app.services.drawio_export import save_drawio_png
                png_url, _ = save_drawio_png(xml)
                if png_url:
                    logger.info(f"[DRAWIO] 已渲染并保存")
                    return f"\n\n![图表]({png_url})\n\n"
                else:
                    # 渲染失败：保留原始 [DRAWIO] 标记让前端 DiagramImage 兜底
                    return match.group(0)
            except Exception as e:
                logger.error(f"[DRAWIO] 处理异常: {e}")
                return match.group(0)

        content = re.sub(r'\[DRAWIO\]([\s\S]*?)\[/DRAWIO\]', _process_drawio_block, content)

    db_message = ChatMessageModel(
        session_id=request.chat_id,
        role=request.role,
        content=content,
        reasoning_content=request.reasoning_content,
        citations=request.citations if request.citations else None,
        model="",
        created_at=datetime.now()
    )
    db.add(db_message)

    db_session = db.query(ChatSessionModel).filter(ChatSessionModel.id == request.chat_id).first()
    if db_session:
        db_session.message_count = db.query(ChatMessageModel).filter(ChatMessageModel.session_id == request.chat_id).count() + 1
        db_session.updated_at = datetime.now()

    db.commit()
    db.refresh(db_message)

    return {
        "id": str(db_message.id),
        "chat_id": db_message.session_id,
        "role": db_message.role,
        "content": db_message.content,
        "created_at": _fmt_iso(db_message.created_at)
    }


class UpdateMessageRequest(BaseModel):
    content: str
    role: Optional[str] = None


@router.put("/messages/{message_id}")
async def update_message(
    message_id: str,
    request: UpdateMessageRequest,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """更新已有消息的内容（用于追加思维导图等场景，避免创建重复消息）"""
    db_message = db.query(ChatMessageModel).filter(
        ChatMessageModel.id == message_id
    ).first()

    if not db_message:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="消息不存在"
        )

    # 验证消息所属的会话属于当前用户
    session = db.query(ChatSessionModel).filter(
        ChatSessionModel.id == db_message.session_id,
        ChatSessionModel.user_id == current_user.student_id
    ).first()
    if not session:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="无权修改此消息"
        )

    content = request.content

    # ═══ 处理 [PLOT] 代码块：执行 → 保存 PNG → 替换为图片 URL ═══
    if request.role == "assistant" and "[PLOT]" in content:
        def _process_plot_block(match):
            code = match.group(1).strip()
            code = re.sub(r'^```[a-zA-Z]*\n?', '', code)
            code = re.sub(r'\n?```$', '', code)
            try:
                from app.api.endpoints.code_execution import _execute_plot
                result = _execute_plot(code)
                if result.success and result.image:
                    b64_data = result.image
                    if b64_data.startswith("data:image/png;base64,"):
                        b64_data = b64_data[len("data:image/png;base64,"):]
                    import base64
                    img_bytes = base64.b64decode(b64_data)
                    plot_id = uuid.uuid4().hex[:12]
                    filename = f"plot_{plot_id}.png"
                    plot_dir = os.path.join("uploads", "plots")
                    os.makedirs(plot_dir, exist_ok=True)
                    filepath = os.path.join(plot_dir, filename)
                    with open(filepath, "wb") as f:
                        f.write(img_bytes)
                    url = f"/api/v1/chat/plots/{filename}"
                    logger.info(f"[PLOT] 已保存: {filepath} ({len(img_bytes)} bytes)")
                    return f"\n\n![图表]({url})\n\n"
                else:
                    logger.warning(f"[PLOT] 执行失败: {result.stderr}")
                    return f"\n\n> ⚠️ 图表生成失败: {result.stderr}\n\n"
            except Exception as e:
                logger.error(f"[PLOT] 处理异常: {e}")
                return f"\n\n> ⚠️ 图表处理异常: {str(e)}\n\n"

        content = re.sub(r'\[PLOT\]([\s\S]*?)\[/PLOT\]', _process_plot_block, content)

    # ═══ 处理 [DRAWIO] 代码块：尝试渲染为 PNG → 替换为图片 URL ═══
    if request.role == "assistant" and "[DRAWIO]" in content:
        def _process_drawio_block(match):
            xml = match.group(1).strip()
            xml = re.sub(r'^```[a-zA-Z]*\n?', '', xml)
            xml = re.sub(r'\n?```$', '', xml)
            if not xml.strip():
                return "\n\n> ⚠️ 图表内容为空\n\n"
            try:
                from app.services.drawio_export import save_drawio_png
                png_url, _ = save_drawio_png(xml)
                if png_url:
                    logger.info(f"[DRAWIO] 已渲染并保存")
                    return f"\n\n![图表]({png_url})\n\n"
                else:
                    return match.group(0)
            except Exception as e:
                logger.error(f"[DRAWIO] 处理异常: {e}")
                return match.group(0)

        content = re.sub(r'\[DRAWIO\]([\s\S]*?)\[/DRAWIO\]', _process_drawio_block, content)

    db_message.content = content

    db_session = db.query(ChatSessionModel).filter(ChatSessionModel.id == db_message.session_id).first()
    if db_session:
        db_session.updated_at = datetime.now()

    db.commit()
    db.refresh(db_message)

    return {
        "id": str(db_message.id),
        "chat_id": db_message.session_id,
        "role": db_message.role,
        "content": db_message.content,
        "created_at": _fmt_iso(db_message.created_at)
    }


@router.post("/attachments")
async def create_attachment(
    request: AttachmentCreate,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    db_attachment = ChatAttachmentModel(
        session_id=request.session_id,
        file_id=request.file_id,
        file_name=request.file_name,
        file_type=request.file_type,
        created_at=datetime.now()
    )
    db.add(db_attachment)
    db.commit()
    db.refresh(db_attachment)

    return {
        "id": str(db_attachment.id),
        "session_id": db_attachment.session_id,
        "file_id": db_attachment.file_id,
        "file_name": db_attachment.file_name,
        "file_type": db_attachment.file_type,
        "created_at": _fmt_iso(db_attachment.created_at)
    }


@router.get("/{chat_id}/attachments")
async def get_attachments(
    chat_id: str,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    attachments = db.query(ChatAttachmentModel).filter(
        ChatAttachmentModel.session_id == chat_id
    ).all()

    return {
        "attachments": [
            {
                "id": str(a.id),
                "session_id": a.session_id,
                "file_id": a.file_id,
                "file_name": a.file_name,
                "file_type": a.file_type,
                "created_at": _fmt_iso(a.created_at)
            }
            for a in attachments
        ]
    }


@router.delete("/attachments/{attachment_id}")
async def delete_attachment(
    attachment_id: str,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    attachment = db.query(ChatAttachmentModel).filter(
        ChatAttachmentModel.id == attachment_id
    ).first()

    if not attachment:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="附件不存在"
        )

    db.delete(attachment)
    db.commit()

    return {"success": True}


# Profile Init API (保留原有功能)
class ChatMessage(BaseModel):
    role: str
    content: str


class ProfileInitRequest(BaseModel):
    message: str
    conversation_history: List[ChatMessage] = []
    collected_info: Dict[str, str] = Field(default_factory=dict)


class ProfileInitResponse(BaseModel):
    reply: str
    collected_info: Dict[str, str]
    is_complete: bool


SYSTEM_PROMPT = """你是一个学习助手，正在帮助用户初始化学习画像。你需要收集以下信息：
1. 专业背景（major）
2. 年级（grade）：大一/大二/大三/大四/研一/研二/研三
3. 学习目标（learningGoal）
4. 学习风格偏好（preferredStyle）：visual（视觉型）/auditory（听觉型）/reading_writing（阅读型）/kinesthetic（实践型）/mixed（混合型）
5. 自我评估基础水平（selfAssessment）

当前已收集的信息：{collected_info}

请根据以下规则进行对话：
- 如果某项信息还未收集，请询问该项信息
- 如果用户回答了，请确认并继续询问下一项
- 当所有5项信息都收集完成后，请回复以"[COMPLETE]"开头的消息
- 保持对话自然流畅，使用友好的语气
- 使用中文回复
"""


def get_collected_info_status(collected_info: Dict[str, str]) -> str:
    required_fields = ['major', 'grade', 'learningGoal', 'preferredStyle', 'selfAssessment']
    status_lines = []
    for field in required_fields:
        value = collected_info.get(field, '')
        if value:
            status_lines.append(f"✓ {field}: {value}")
        else:
            status_lines.append(f"✗ {field}: 未收集")
    return "\n".join(status_lines)


async def call_deepseek_api(messages: List[Dict[str, str]]) -> str:
    api_key = settings.DEEPSEEK_API_KEY
    model = settings.DEEPSEEK_MODEL
    base_url = settings.DEEPSEEK_BASE_URL

    if not api_key:
        api_key = settings.QWEN_API_KEY
        base_url = settings.QWEN_BASE_URL
        model = settings.QWEN_MODEL

    if not api_key:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="AI 服务未配置，请在「API 设置」页面配置 DeepSeek 或 Qwen API Key 后再使用智能对话功能"
        )

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    payload = {
        "model": model,
        "messages": messages,
        "temperature": settings.DEEPSEEK_TEMPERATURE,
        "max_tokens": settings.DEEPSEEK_MAX_TOKENS,
    }

    async with httpx.AsyncClient(timeout=settings.DEEPSEEK_TIMEOUT) as client:
        response = await client.post(
            f"{base_url}/chat/completions",
            headers=headers,
            json=payload
        )

        if response.status_code != 200:
            raise HTTPException(
                status_code=status.HTTP_502_BAD_GATEWAY,
                detail=f"DeepSeek API 调用失败: {response.text}"
            )

        result = response.json()
        return result["choices"][0]["message"]["content"]


@router.post("/profile-init", response_model=ProfileInitResponse)
async def profile_init_chat(
    request: ProfileInitRequest,
    current_user: CurrentUser = Depends(get_current_user)
):
    try:
        collected_info = request.collected_info.copy()

        required_fields = ['major', 'grade', 'learningGoal', 'preferredStyle', 'selfAssessment']

        user_message = request.message.lower()

        if 'major' not in collected_info or not collected_info['major']:
            collected_info['major'] = request.message
        elif 'grade' not in collected_info or not collected_info['grade']:
            collected_info['grade'] = request.message
        elif 'learningGoal' not in collected_info or not collected_info['learningGoal']:
            collected_info['learningGoal'] = request.message
        elif 'preferredStyle' not in collected_info or not collected_info['preferredStyle']:
            style_map = {
                '视觉': 'visual',
                '听觉': 'auditory',
                '阅读': 'reading_writing',
                '实践': 'kinesthetic',
                'visual': 'visual',
                'auditory': 'auditory',
                'reading': 'reading_writing',
                'kinesthetic': 'kinesthetic',
                'mixed': 'mixed',
                '混合': 'mixed',
            }
            for keyword, style in style_map.items():
                if keyword in user_message:
                    collected_info['preferredStyle'] = style
                    break
            else:
                collected_info['preferredStyle'] = request.message
        elif 'selfAssessment' not in collected_info or not collected_info['selfAssessment']:
            collected_info['selfAssessment'] = request.message

        all_collected = all(collected_info.get(field) for field in required_fields)

        info_status = get_collected_info_status(collected_info)
        system_prompt = SYSTEM_PROMPT.format(collected_info=info_status)

        messages = [{"role": "system", "content": system_prompt}]

        for msg in request.conversation_history[-10:]:
            messages.append({"role": msg.role, "content": msg.content})

        messages.append({"role": "user", "content": request.message})

        llm_reply = await call_deepseek_api(messages)

        is_complete = "[COMPLETE]" in llm_reply.upper()
        reply = llm_reply.replace("[COMPLETE]", "").strip()

        if not reply:
            reply = "很好！画像初始化完成。"
            is_complete = True

        return ProfileInitResponse(
            reply=reply,
            collected_info=collected_info,
            is_complete=is_complete
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"对话处理失败: {str(e)}"
        )


# ═══ 图表静态文件服务 ═══
@router.get("/plots/{filename}")
async def serve_plot_image(filename: str):
    """提供保存的 matplotlib 图表 PNG 文件"""
    filepath = os.path.join("uploads", "plots", filename)
    if not os.path.isfile(filepath):
        raise HTTPException(status_code=404, detail="图表文件不存在")
    return FileResponse(filepath, media_type="image/png")


@router.get("/drawio/{filename}")
async def serve_drawio_image(filename: str):
    """提供保存的 draw.io 图表 PNG 文件"""
    filepath = os.path.join("uploads", "drawio", filename)
    if not os.path.isfile(filepath):
        raise HTTPException(status_code=404, detail="图表文件不存在")
    return FileResponse(filepath, media_type="image/png")


@router.get("/drawio/{drawio_id}/xml")
async def serve_drawio_xml(drawio_id: str):
    """提供保存的 draw.io 原始 XML 文件（用于编辑器打开）"""
    # drawio_id is e.g. "drawio_abc123" from filename "drawio_abc123.png"
    clean_id = drawio_id.replace("drawio_", "")
    filepath = os.path.join("uploads", "drawio", f"drawio_{clean_id}.xml")
    if not os.path.isfile(filepath):
        raise HTTPException(status_code=404, detail="XML 文件不存在")
    return FileResponse(filepath, media_type="application/xml")


# ═══ 预测下次提问 ═══

class NextQuestionsRequest(BaseModel):
    conversation_history: List[Dict[str, str]]


class NextQuestionsResponse(BaseModel):
    questions: List[str]


@router.post("/next-questions", response_model=NextQuestionsResponse)
async def predict_next_questions(
    request: NextQuestionsRequest,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """根据对话历史预测用户最可能追问的 2-3 个问题"""
    from app.crud.api_settings import api_settings_crud

    # 获取用户 API Key
    api_info = None
    for provider in ["qwen", "deepseek"]:
        api = api_settings_crud.get_setting_value(db, str(current_user.student_id), provider)
        if api:
            api_info = api
            break

    if not api_info:
        return NextQuestionsResponse(questions=[])

    base_url = api_info.get("base_url") or "https://dashscope.aliyuncs.com/compatible-mode/v1"
    model = api_info.get("model_version") or "qwen-turbo"
    api_key = api_info.get("api_key")

    if not api_key:
        return NextQuestionsResponse(questions=[])

    # 取最近 3 轮对话（6 条消息）
    recent = request.conversation_history[-6:]

    # 构建提示词上下文
    context_lines = []
    for msg in recent:
        role_label = "用户" if msg.get("role") == "user" else "助手"
        content = msg.get("content", "")
        context_lines.append(f"{role_label}: {content}")
    context = "\n".join(context_lines)

    prompt = (
        "你是一个学习助手。根据以下对话历史，预测用户接下来最可能追问的 2-3 个问题。\n"
        "每个问题不超过 15 个字。只返回 JSON 格式：\n"
        '{"questions": ["问题1", "问题2", "问题3"]}\n\n'
        f"对话历史：\n{context}"
    )

    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            response = await client.post(
                f"{base_url}/chat/completions",
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": model,
                    "messages": [{"role": "user", "content": prompt}],
                    "temperature": 0.3,
                },
            )

            if response.status_code != 200:
                logger.warning(f"next-questions LLM 调用失败: status={response.status_code}")
                return NextQuestionsResponse(questions=[])

            result = response.json()
            content = result["choices"][0]["message"]["content"].strip()

            # 尝试直接解析 JSON
            parsed = None
            try:
                parsed = json.loads(content)
            except json.JSONDecodeError:
                pass

            # 尝试从 markdown 代码块提取
            if parsed is None:
                m = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', content, re.DOTALL)
                if m:
                    try:
                        parsed = json.loads(m.group(1))
                    except json.JSONDecodeError:
                        pass

            # 尝试在文本中查找 JSON 对象
            if parsed is None:
                m = re.search(r'\{(?:[^{}]|\{[^{}]*\})*"questions"\s*:\s*\[[^\]]*\](?:[^{}]|\{[^{}]*\})*\}', content)
                if m:
                    try:
                        parsed = json.loads(m.group(0))
                    except json.JSONDecodeError:
                        pass

            if parsed is None or not isinstance(parsed, dict):
                logger.warning(f"next-questions: 无法解析 LLM 响应为 JSON: {content[:200]}")
                return NextQuestionsResponse(questions=[])

            questions = parsed.get("questions", [])
            if not isinstance(questions, list):
                questions = []

            # 截断每个问题到 15 字，最多 3 个
            questions = [q[:15] for q in questions[:3]]
            return NextQuestionsResponse(questions=questions)

    except httpx.TimeoutException:
        logger.warning("next-questions: LLM 调用超时")
        return NextQuestionsResponse(questions=[])
    except Exception as e:
        logger.warning(f"next-questions: 异常: {e}")
        return NextQuestionsResponse(questions=[])
