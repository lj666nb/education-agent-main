import os
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Query
from sqlalchemy.orm import Session
from typing import List

from app.db.database import get_db
from app.api.dependencies import get_current_user
from app.crud import project as project_crud
from app.schemas.project import (
    ProjectCreate, ProjectUpdate, ProjectResponse,
    PromptCreate, PromptUpdate, PromptResponse,
    DocumentCreate, DocumentResponse, RetrievalRequest, RetrievalResponse
)

router = APIRouter(prefix="/projects", tags=["projects"])


def extract_pdf_text(raw_content: bytes, max_chars: int = 50000) -> str:
    try:
        import fitz
        doc = fitz.open(stream=raw_content, filetype="pdf")
        total_pages = len(doc)
        page_texts = []
        for page_num in range(total_pages):
            page = doc[page_num]
            page_text = page.get_text().strip()
            if page_text:
                page_texts.append(f"[第{page_num + 1}页] {page_text}")
        doc.close()

        if page_texts:
            full_text = "\n\n".join(page_texts)
            if len(full_text) > max_chars:
                full_text = full_text[:max_chars] + f"\n\n... [内容已截断，仅显示前{max_chars}字符]"
            return full_text
        elif total_pages > 0:
            return f"[该PDF共{total_pages}页，未能提取到可识别的文本内容。可能原因：该PDF为扫描件或图片型PDF，文字以图像形式存储。建议使用OCR工具提取文字后重新输入。]"
        else:
            return "[PDF文件为空]"
    except ImportError:
        return "[PDF解析库未安装，无法提取文件内容，请安装PyMuPDF]"
    except Exception as e:
        return f"[PDF文件解析失败: {str(e)}]"


def extract_pptx_text(raw_content: bytes, max_chars: int = 50000) -> str:
    try:
        from pptx import Presentation
        from io import BytesIO
        prs = Presentation(BytesIO(raw_content))
        total_slides = len(prs.slides)
        slide_texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(row_texts))
            if texts:
                slide_texts.append(f"[第{slide_num}页] " + "\n".join(texts))
        if slide_texts:
            full_text = "\n\n".join(slide_texts)
            if len(full_text) > max_chars:
                full_text = full_text[:max_chars] + f"\n\n... [内容已截断，仅显示前{max_chars}字符]"
            return full_text
        elif total_slides > 0:
            return f"[该PPT共{total_slides}页，未能提取到可识别的文本内容]"
        else:
            return "[PPT文件为空]"
    except ImportError:
        return "[PPT解析库未安装，无法提取文件内容，请安装python-pptx]"
    except Exception as e:
        return f"[PPT文件解析失败: {str(e)}]"


def extract_docx_text(raw_content: bytes, max_chars: int = 50000) -> str:
    try:
        from docx import Document
        from io import BytesIO
        doc = Document(BytesIO(raw_content))
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())
        if paragraphs:
            full_text = "\n\n".join(paragraphs)
            if len(full_text) > max_chars:
                full_text = full_text[:max_chars] + f"\n\n... [内容已截断，仅显示前{max_chars}字符]"
            return full_text
        return "[Word文件为空]"
    except ImportError:
        return "[Word解析库未安装，无法提取文件内容，请安装python-docx]"
    except Exception as e:
        return f"[Word文件解析失败: {str(e)}]"


@router.post("/", response_model=ProjectResponse)
async def create_project(
    project: ProjectCreate,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    return project_crud.create_project(db, current_user.student_id, project)


@router.get("/", response_model=List[ProjectResponse])
async def get_projects(
    skip: int = 0,
    limit: int = 100,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    projects = project_crud.get_projects(db, current_user.student_id, skip, limit)
    return [
        ProjectResponse(
            id=p.id,
            user_id=p.user_id,
            name=p.name,
            description=p.description,
            created_at=p.created_at,
            updated_at=p.updated_at,
            prompt_count=len(p.prompts) if p.prompts else 0,
            document_count=len(p.documents) if p.documents else 0
        )
        for p in projects
    ]


@router.get("/{project_id}", response_model=ProjectResponse)
async def get_project(
    project_id: str,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = project_crud.get_project(db, project_id, current_user.student_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")
    return ProjectResponse(
        id=project.id,
        user_id=project.user_id,
        name=project.name,
        description=project.description,
        created_at=project.created_at,
        updated_at=project.updated_at,
        prompt_count=len(project.prompts) if project.prompts else 0,
        document_count=len(project.documents) if project.documents else 0
    )


@router.put("/{project_id}", response_model=ProjectResponse)
async def update_project(
    project_id: str,
    project: ProjectUpdate,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    db_project = project_crud.update_project(db, project_id, current_user.student_id, project)
    if not db_project:
        raise HTTPException(status_code=404, detail="项目不存在")
    return ProjectResponse(
        id=db_project.id,
        user_id=db_project.user_id,
        name=db_project.name,
        description=db_project.description,
        created_at=db_project.created_at,
        updated_at=db_project.updated_at,
        prompt_count=len(db_project.prompts) if db_project.prompts else 0,
        document_count=len(db_project.documents) if db_project.documents else 0
    )


@router.delete("/{project_id}")
async def delete_project(
    project_id: str,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    success = project_crud.delete_project(db, project_id, current_user.student_id)
    if not success:
        raise HTTPException(status_code=404, detail="项目不存在")
    return {"message": "项目已删除"}


@router.post("/{project_id}/prompts", response_model=PromptResponse)
async def create_prompt(
    project_id: str,
    prompt: PromptCreate,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = project_crud.get_project(db, project_id, current_user.student_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")
    return project_crud.create_prompt(db, project_id, prompt)


@router.get("/{project_id}/prompts", response_model=List[PromptResponse])
async def get_prompts(
    project_id: str,
    active_only: bool = False,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = project_crud.get_project(db, project_id, current_user.student_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")
    return project_crud.get_prompts(db, project_id, active_only)


@router.put("/{project_id}/prompts/{prompt_id}", response_model=PromptResponse)
async def update_prompt(
    project_id: str,
    prompt_id: str,
    prompt: PromptUpdate,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    db_prompt = project_crud.update_prompt(db, prompt_id, prompt)
    if not db_prompt:
        raise HTTPException(status_code=404, detail="提示词不存在")
    return db_prompt


@router.delete("/{project_id}/prompts/{prompt_id}")
async def delete_prompt(
    project_id: str,
    prompt_id: str,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    success = project_crud.delete_prompt(db, prompt_id)
    if not success:
        raise HTTPException(status_code=404, detail="提示词不存在")
    return {"message": "提示词已删除"}


@router.post("/{project_id}/documents", response_model=DocumentResponse)
async def create_document(
    project_id: str,
    doc: DocumentCreate,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = project_crud.get_project(db, project_id, current_user.student_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")

    document = project_crud.create_document(db, project_id, doc)

    if doc.content_text:
        chunk_size = 500
        content = doc.content_text
        chunks = []
        for i in range(0, len(content), chunk_size):
            chunks.append({"content": content[i:i+chunk_size]})
        if chunks:
            project_crud.create_chunks(db, document.id, chunks)

    return document


@router.post("/{project_id}/documents/upload", response_model=DocumentResponse)
async def upload_document(
    project_id: str,
    file: UploadFile = File(...),
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = project_crud.get_project(db, project_id, current_user.student_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")

    raw_content = await file.read()
    MAX_FILE_SIZE = 50 * 1024 * 1024
    if len(raw_content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="文件过大，单个文件不能超过 50MB，请压缩后重新上传")
    file_ext = os.path.splitext(file.filename or "")[1].lower() if file.filename else ""
    is_pdf = file_ext == ".pdf" or (file.content_type and file.content_type == "application/pdf")
    is_pptx = file_ext == ".pptx" or (file.content_type and file.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    is_ppt = file_ext == ".ppt" or (file.content_type and file.content_type == "application/vnd.ms-powerpoint")
    is_docx = file_ext == ".docx" or (file.content_type and file.content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    is_doc = file_ext == ".doc" or (file.content_type and file.content_type == "application/msword")

    if is_pdf:
        content_text = extract_pdf_text(raw_content)
    elif is_pptx:
        content_text = extract_pptx_text(raw_content)
    elif is_ppt:
        content_text = "[旧版PPT格式（.ppt），无法直接解析文本内容。建议将文件另存为 .pptx 格式后重新上传。]"
    elif is_docx:
        content_text = extract_docx_text(raw_content)
    elif is_doc:
        content_text = "[旧版Word格式（.doc），无法直接解析文本内容。建议将文件另存为 .docx 格式后重新上传。]"
    else:
        content_text = raw_content.decode('utf-8', errors='ignore').replace('\x00', '')

    doc_data = DocumentCreate(
        name=file.filename or '未命名文件',
        file_type=file.content_type or 'text/plain',
        file_size=len(raw_content),
        content_text=content_text
    )

    document = project_crud.create_document(db, project_id, doc_data)

    if content_text:
        chunk_size = 500
        chunks = []
        for i in range(0, len(content_text), chunk_size):
            chunks.append({"content": content_text[i:i+chunk_size]})
        if chunks:
            project_crud.create_chunks(db, document.id, chunks)

    return document


@router.get("/{project_id}/documents", response_model=List[DocumentResponse])
async def get_documents(
    project_id: str,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = project_crud.get_project(db, project_id, current_user.student_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")
    return project_crud.get_documents(db, project_id)


@router.delete("/{project_id}/documents/{document_id}")
async def delete_document(
    project_id: str,
    document_id: str,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    success = project_crud.delete_document(db, document_id)
    if not success:
        raise HTTPException(status_code=404, detail="文档不存在")
    return {"message": "文档已删除"}


@router.post("/{project_id}/retrieve", response_model=RetrievalResponse)
async def retrieve_documents(
    project_id: str,
    request: RetrievalRequest,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = project_crud.get_project(db, project_id, current_user.student_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")

    chunks = project_crud.retrieve_relevant_chunks(
        db, project_id, request.query, request.top_k,
        user_id=current_user.student_id
    )
    return RetrievalResponse(chunks=chunks, query=request.query)


@router.get("/{project_id}/build-index/estimate")
async def estimate_build_index(
    project_id: str,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = project_crud.get_project(db, project_id, current_user.student_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")

    chunk_count = project_crud.count_project_chunks(db, project_id)
    if chunk_count == 0:
        return {
            "can_build": False,
            "message": "项目中没有任何文档块，请先上传文档",
            "chunk_count": 0
        }

    from app.core.vector_store import estimate_build_time
    estimate = estimate_build_time(chunk_count)
    estimate["can_build"] = True
    estimate["project_id"] = project_id
    return estimate


@router.post("/{project_id}/build-index")
async def build_index(
    project_id: str,
    background: bool = Query(True, description="是否在后台构建"),
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = project_crud.get_project(db, project_id, current_user.student_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")

    api_key_info = project_crud.get_text_embedding_api_key(db, current_user.student_id)
    if not api_key_info:
        raise HTTPException(
            status_code=400,
            detail="请先在API设置中配置text_embedding的API密钥"
        )

    chunk_count = project_crud.count_project_chunks(db, project_id)
    if chunk_count == 0:
        raise HTTPException(
            status_code=400,
            detail="项目中没有任何文档块，请先上传文档"
        )

    from app.core.vector_store import estimate_build_time, get_build_status
    estimate = estimate_build_time(chunk_count)

    existing_status = get_build_status(project_id)
    if existing_status and existing_status["status"] == "building":
        raise HTTPException(
            status_code=409,
            detail="该项目正在构建向量索引中，请等待当前任务完成"
        )

    if background:
        message = project_crud.build_vector_index(db, project_id, current_user.student_id)
        return {
            "background": True,
            "message": message,
            "project_id": project_id,
            "chunk_count": chunk_count,
            "estimated_time": estimate["estimated_time"],
            "status_endpoint": f"/api/v1/projects/{project_id}/build-index/status"
        }
    else:
        from app.core.vector_store import build_project_index_sync
        from app.db.database import SessionLocal
        build_project_index_sync(
            SessionLocal, project_id, str(current_user.student_id),
            api_key_info["api_key"],
            model_name=api_key_info.get("model_version")
        )
        return {
            "background": False,
            "message": "向量索引构建完成",
            "project_id": project_id,
            "chunk_count": chunk_count
        }


@router.get("/{project_id}/build-index/status")
async def get_build_index_status(
    project_id: str,
    current_user=Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = project_crud.get_project(db, project_id, current_user.student_id)
    if not project:
        raise HTTPException(status_code=404, detail="项目不存在")

    from app.core.vector_store import get_build_status
    status = get_build_status(project_id)
    if not status:
        from app.core.vector_store import get_vector_store
        store = get_vector_store(project_id)
        if store.ntotal > 0:
            return {
                "status": "completed",
                "message": f"向量索引已就绪，共 {store.ntotal} 个文档块"
            }
        return {
            "status": "not_started",
            "message": "尚未开始构建向量索引"
        }

    return status